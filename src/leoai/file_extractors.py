from __future__ import annotations

from io import BytesIO
import json
import zipfile


def _extract_pdf_text(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text.strip())
    return "\n\n".join(parts).strip()


def _extract_docx_text(data: bytes) -> str:
    from docx import Document

    doc = Document(BytesIO(data))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n".join(paragraphs).strip()


def _extract_xlsx_text(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(filename=BytesIO(data), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"[Sheet] {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines).strip()


def _extract_pptx_text(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(BytesIO(data))
    lines: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"[Slide {idx}]")
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            text = text.strip()
            if text:
                lines.append(text)
    return "\n".join(lines).strip()


def _extract_image_metadata(data: bytes) -> str:
    from PIL import Image

    with Image.open(BytesIO(data)) as img:
        meta = {
            "format": img.format,
            "mode": img.mode,
            "width": img.width,
            "height": img.height,
        }
        return json.dumps(meta, ensure_ascii=False)


def _extract_text_like(data: bytes) -> str:
    return data.decode("utf-8", errors="replace").strip()


def extract_content_for_rag(filename: str, content_type: str, data: bytes) -> tuple[str, dict[str, str]]:
    lowered_name = (filename or "").lower()
    lowered_type = (content_type or "").lower()
    meta: dict[str, str] = {"strategy": "binary_metadata"}

    try:
        if lowered_name.endswith(".pdf") or lowered_type == "application/pdf":
            text = _extract_pdf_text(data)
            meta["strategy"] = "pdf_text"
            return text or "", meta
        if lowered_name.endswith(".docx") or lowered_type in {
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        }:
            text = _extract_docx_text(data)
            meta["strategy"] = "docx_text"
            return text or "", meta
        if lowered_name.endswith((".xlsx", ".xlsm")) or lowered_type in {
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel.sheet.macroenabled.12",
            "application/vnd.ms-excel",
        }:
            text = _extract_xlsx_text(data)
            meta["strategy"] = "xlsx_text"
            return text or "", meta
        if lowered_name.endswith(".pptx") or lowered_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text = _extract_pptx_text(data)
            meta["strategy"] = "pptx_text"
            return text or "", meta
        if lowered_type.startswith("image/") or lowered_name.endswith((".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff")):
            text = _extract_image_metadata(data)
            meta["strategy"] = "image_metadata"
            return text, meta
        if lowered_type.startswith("text/") or lowered_name.endswith(
            (".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml", ".log", ".py", ".js", ".ts", ".html", ".css")
        ):
            text = _extract_text_like(data)
            meta["strategy"] = "plain_text"
            return text, meta
        if zipfile.is_zipfile(BytesIO(data)):
            meta["strategy"] = "zip_metadata"
            return f"Arquivo ZIP recebido: {filename} ({len(data)} bytes)", meta
    except Exception as exc:
        meta["strategy"] = "fallback_after_error"
        meta["error"] = str(exc)[:200]

    return (
        f"Arquivo binario enviado: {filename}\n"
        f"Content-Type: {content_type or 'application/octet-stream'}\n"
        f"Tamanho: {len(data)} bytes",
        meta,
    )
